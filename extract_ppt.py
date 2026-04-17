from pptx import Presentation
from pptx.util import Inches, Emu
import json

prs = Presentation(r'C:\Users\ZDominguez\OneDrive - arthrex.com\Desktop\robot puzzles.pptx')

for i, slide in enumerate(prs.slides):
    title = ''
    for shape in slide.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if txt:
                title = txt
                break
    print(f'\n=== Slide {i+1}: {title} ===')
    for shape in slide.shapes:
        name = shape.name
        l = shape.left
        t = shape.top
        w = shape.width
        h = shape.height
        fill_color = ''
        try:
            if hasattr(shape, 'fill') and shape.fill and shape.fill.fore_color:
                fill_color = str(shape.fill.fore_color.rgb)
        except:
            pass
        txt = ''
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
        print(f'  {name}: pos=({l},{t}) size=({w},{h}) fill={fill_color} text="{txt}"')
