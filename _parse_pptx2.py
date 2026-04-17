from pptx import Presentation
from pptx.util import Emu
import json, math

prs = Presentation(r'C:\Users\ZDominguez\OneDrive - arthrex.com\Desktop\robot puzzles.pptx')

for i, slide in enumerate(prs.slides):
    print(f'=== SLIDE {i+1} ===')
    shapes = []
    for shape in slide.shapes:
        txt = ''
        if shape.has_text_frame:
            txt = ' '.join(p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip())
        # Convert EMU to inches for readability
        l = round(shape.left / 914400, 2)
        t = round(shape.top / 914400, 2)
        w = round(shape.width / 914400, 2)
        h = round(shape.height / 914400, 2)
        cx = round(l + w/2, 2)
        cy = round(t + h/2, 2)
        shapes.append({'name': shape.name, 'text': txt, 'left': l, 'top': t, 'w': w, 'h': h, 'cx': cx, 'cy': cy})
        label = txt if txt else '(empty)'
        print(f'  {label:20s} cx={cx:6.2f} cy={cy:6.2f}  [{l},{t}]-[{round(l+w,2)},{round(t+h,2)}]  size={w}x{h}')
    print()
