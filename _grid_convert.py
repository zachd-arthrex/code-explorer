from pptx import Presentation
import math, json

prs = Presentation(r'C:\Users\ZDominguez\OneDrive - arthrex.com\Desktop\robot puzzles.pptx')

for i, slide in enumerate(prs.slides):
    shapes = []
    for shape in slide.shapes:
        txt = ''
        if shape.has_text_frame:
            txt = ' '.join(p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()).lower()
        cx = (shape.left + shape.width/2) / 914400
        cy = (shape.top + shape.height/2) / 914400
        w = shape.width / 914400
        shapes.append({'cx': cx, 'cy': cy, 'w': w, 'text': txt})

    if not shapes:
        continue

    # Determine cell size from shape widths (most common)
    widths = [s['w'] for s in shapes]
    cell = round(max(set(widths), key=widths.count), 2)

    # Quantize to grid coords
    all_cx = sorted(set(round(s['cx'], 1) for s in shapes))
    all_cy = sorted(set(round(s['cy'], 1) for s in shapes))

    # Cluster nearby coordinates (within cell/3)
    def cluster(vals, threshold):
        clusters = []
        for v in sorted(vals):
            if clusters and abs(v - clusters[-1][-1]) < threshold:
                clusters[-1].append(v)
            else:
                clusters.append([v])
        return [sum(c)/len(c) for c in clusters]

    cx_clusters = cluster([s['cx'] for s in shapes], cell * 0.4)
    cy_clusters = cluster([s['cy'] for s in shapes], cell * 0.4)

    def nearest(val, clusters):
        return min(range(len(clusters)), key=lambda i: abs(clusters[i] - val))

    # Build grid
    rows = len(cy_clusters)
    cols = len(cx_clusters)
    grid = [[-1]*cols for _ in range(rows)]

    robot = None
    goal = None
    collectibles = []
    portals = {}

    for s in shapes:
        c = nearest(s['cx'], cx_clusters)
        r = nearest(s['cy'], cy_clusters)
        txt = s['text']

        if 'robot' in txt:
            robot = [r, c]
            grid[r][c] = 0
        elif 'goal' in txt:
            goal = [r, c]
            grid[r][c] = 2
        elif 'power' in txt:
            collectibles.append([r, c])
            grid[r][c] = 0  # collectible stored separately
        elif 'portal' in txt:
            # Extract portal label (A, B, C, or just "portal")
            label = 'A'  # default for unlabeled portals
            for p in ['a', 'b', 'c']:
                if p in txt.replace('portal', '').strip():
                    label = p.upper()
                    break
            if label not in portals:
                portals[label] = []
            portals[label].append([r, c])
            grid[r][c] = 0
        else:
            # plain empty cell
            if grid[r][c] == -1:
                grid[r][c] = 0

    portal_pairs = []
    for label in sorted(portals.keys()):
        coords = portals[label]
        if len(coords) == 2:
            portal_pairs.append(coords)

    print(f'// === Level {i+1} === ({rows}x{cols}, cell={cell})')
    print(f'// Robot: {robot}, Goal: {goal}')
    if collectibles:
        print(f'// Collectibles: {collectibles}')
    if portal_pairs:
        print(f'// Portals: {portal_pairs}')
    print('{')
    print(f'  id: {i+1}, title: "Level {i+1}",')
    print(f'  desc: "",')
    print(f'  grid: [')
    for row in grid:
        print(f'    [{",".join(str(v).rjust(2) for v in row)}],')
    print(f'  ],')
    print(f'  player: {robot}, startDir: 1, cols: {cols}, rows: {rows},')
    print(f'  functionsUnlocked: false, minSteps: null, maxSteps: null,')
    print(f'  collectibles: {collectibles},')
    if portal_pairs:
        print(f'  portals: {json.dumps(portal_pairs)},')
    print('},')
    print()
