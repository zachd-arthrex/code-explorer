from pptx import Presentation
import math

prs = Presentation(r'C:\Users\ZDominguez\OneDrive - arthrex.com\Desktop\robot puzzles.pptx')

for i, slide in enumerate(prs.slides):
    # Collect all rectangles (not text boxes)
    rects = []
    for shape in slide.shapes:
        if 'Rectangle' in shape.name:
            txt = ''
            if shape.has_text_frame:
                txt = shape.text_frame.text.strip()
            rects.append({
                'name': shape.name,
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height,
                'text': txt
            })
    
    if not rects:
        continue
    
    # Find cell size from first rect
    cell_w = rects[0]['width']
    cell_h = rects[0]['height']
    
    # Snap positions to grid
    min_left = min(r['left'] for r in rects)
    min_top = min(r['top'] for r in rects)
    
    grid_cells = []
    for r in rects:
        col = round((r['left'] - min_left) / cell_w)
        row = round((r['top'] - min_top) / cell_h)
        grid_cells.append((row, col, r['text']))
    
    max_row = max(c[0] for c in grid_cells)
    max_col = max(c[1] for c in grid_cells)
    
    # Build grid: -1 = void, 0 = floor
    grid = [[-1] * (max_col + 1) for _ in range(max_row + 1)]
    
    player = None
    goal = None
    collectibles = []
    portals = {}  # name -> list of (row, col)
    
    for row, col, txt in grid_cells:
        txt_lower = txt.lower().strip()
        grid[row][col] = 0  # floor
        
        if 'robot' in txt_lower:
            player = [row, col]
        elif 'goal' in txt_lower:
            goal = [row, col]
            grid[row][col] = 2
        elif 'power' in txt_lower:
            collectibles.append([row, col])
        elif 'portal' in txt_lower:
            # Extract portal name (A, B, C, or generic)
            parts = txt_lower.replace('portal', '').strip()
            pname = parts if parts else 'X'
            if pname not in portals:
                portals[pname] = []
            portals[pname].append([row, col])
    
    # Format portal pairs
    portal_pairs = []
    for pname, positions in portals.items():
        if len(positions) == 2:
            portal_pairs.append([positions[0], positions[1]])
        elif len(positions) > 2:
            # Pair them up in order
            for j in range(0, len(positions) - 1, 2):
                portal_pairs.append([positions[j], positions[j+1]])
    
    print(f'\n=== Slide {i+1} (Level {i+1}) ===')
    print(f'Grid ({max_row+1}x{max_col+1}):')
    for row in grid:
        print(f'  {row},')
    print(f'Player: {player}  Goal: {goal}')
    print(f'Collectibles: {collectibles}')
    if portal_pairs:
        print(f'Portal pairs: {portal_pairs}')
    print(f'Cell size: {cell_w}x{cell_h}')
