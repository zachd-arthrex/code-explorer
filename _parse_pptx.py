from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.dml.color import RGBColor
import json

prs = Presentation(r'C:\Users\ZDominguez\OneDrive - arthrex.com\Desktop\robot puzzles.pptx')
for i, slide in enumerate(prs.slides):
    print(f'=== SLIDE {i+1} ===')
    for shape in slide.shapes:
        print(f'  Shape: type={shape.shape_type}, name={shape.name}')
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    print(f'    TEXT: {text}')
        if shape.has_table:
            tbl = shape.table
            print(f'    TABLE: {len(list(tbl.rows))} rows x {len(tbl.columns)} cols')
            for ri, row in enumerate(tbl.rows):
                cells = []
                for ci, cell in enumerate(row.cells):
                    txt = cell.text.strip()
                    # Try to get cell fill color
                    fill = cell.fill
                    color_info = ''
                    try:
                        if fill.type is not None:
                            fg = fill.fore_color
                            if fg.type is not None:
                                color_info = f' bg={fg.rgb}'
                    except:
                        pass
                    cells.append(f'{txt}{color_info}')
                print(f'      Row {ri}: {cells}')
    print()
